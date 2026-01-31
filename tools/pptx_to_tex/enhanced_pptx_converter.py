#!/usr/bin/env python3

# Copyright (c) 2026 Slide Forge Team
#
# Permission is hereby granted, free of charge, to any person obtaining a copy
# of this software and associated documentation files (the "Software"), to deal
# in the Software without restriction, including without limitation the rights
# to use, copy, modify, merge, publish, distribute, sublicense, and/or sell
# copies of the Software, and to permit persons to whom the Software is
# furnished to do so, subject to the following conditions:
#
# The above copyright notice and this permission notice shall be included in all
# copies or substantial portions of the Software.
#
# THE SOFTWARE IS PROVIDED "AS IS", WITHOUT WARRANTY OF ANY KIND, EXPRESS OR
# IMPLIED, INCLUDING BUT NOT LIMITED TO THE WARRANTIES OF MERCHANTABILITY,
# FITNESS FOR A PARTICULAR PURPOSE AND NONINFRINGEMENT. IN NO EVENT SHALL THE
# AUTHORS OR COPYRIGHT HOLDERS BE LIABLE FOR ANY CLAIM, DAMAGES OR OTHER
# LIABILITY, WHETHER IN AN ACTION OF CONTRACT, TORT OR OTHERWISE, ARISING FROM,
# OUT OF OR IN CONNECTION WITH THE SOFTWARE OR THE USE OR OTHER DEALINGS IN THE
# SOFTWARE.

"""
Enhanced LaTeX to PowerPoint converter using structured parsing
"""

from advanced_latex_parser import LaTeXBeamerParser
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import logging

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class EnhancedPPTXConverter:
    """Convert structured LaTeX data to high-quality PowerPoint"""

    def __init__(self):
        self.parser = LaTeXBeamerParser()

    def convert_latex_to_pptx(self, latex_file: str, output_file: str):
        """Convert LaTeX file to PowerPoint with enhanced formatting"""

        # Parse LaTeX structure
        logger.info("Parsing LaTeX file...")
        latex_data = self.parser.parse_file(latex_file)

        # Create PowerPoint presentation
        prs = Presentation()

        # Add title slide
        self._create_title_slide(prs, latex_data['metadata'])

        # Convert each frame
        for frame in latex_data['frames']:
            self._convert_frame(prs, frame)

        # Save presentation
        prs.save(output_file)
        logger.info(f"Successfully created {output_file}")

    def _create_title_slide(self, prs: Presentation, metadata: dict):
        """Create title slide from LaTeX metadata"""
        title_layout = prs.slide_layouts[0]  # Title slide layout
        title_slide = prs.slides.add_slide(title_layout)

        # Set title
        title = metadata.get('title', 'Presentation')
        if title_slide.shapes.title:
            title_slide.shapes.title.text = title

        # Set subtitle/author
        subtitle = metadata.get('author', '')
        if subtitle and len(title_slide.placeholders) > 1:
            title_slide.placeholders[1].text = subtitle

    def _convert_frame(self, prs: Presentation, frame: dict):
        """Convert a single frame to PowerPoint slide"""

        # Choose appropriate layout
        if self._has_complex_content(frame):
            slide_layout = prs.slide_layouts[1]  # Title and content
        else:
            slide_layout = prs.slide_layouts[2]  # Section header

        slide = prs.slides.add_slide(slide_layout)

        # Set slide title
        if slide.shapes.title:
            slide.shapes.title.text = frame['title']

        # Process content
        if len(slide.placeholders) > 1:
            content_placeholder = slide.placeholders[1]
            self._add_content_to_placeholder(content_placeholder, frame['elements'])

        # Add images separately for better positioning
        self._add_images_to_slide(slide, frame['elements'])

    def _has_complex_content(self, frame: dict) -> bool:
        """Check if frame has complex content needing full layout"""
        for element in frame['elements']:
            if element['type'] in ['itemize', 'block', 'columns', 'image']:
                return True
        return False

    def _add_content_to_placeholder(self, placeholder, elements: list):
        """Add content to slide placeholder with proper formatting"""
        text_frame = placeholder.text_frame
        text_frame.clear()  # Remove existing text

        # Debug: print element structure
        print(f"Processing {len(elements)} elements:")
        for i, element in enumerate(elements):
            print(f"  {i}: {element['type']} - {element}")

        for element in elements:
            if element['type'] == 'text':
                p = text_frame.add_paragraph()
                p.text = element['text']
                p.font.size = Pt(18)

                # Apply formatting if specified
                if element.get('formatting') == 'bold':
                    p.font.bold = True
                elif element.get('formatting') == 'italic':
                    p.font.italic = True

            elif element['type'] == 'itemize':
                # Handle itemize environment
                for item in element.get('content', []):
                    if item['type'] == 'item':
                        p = text_frame.add_paragraph()
                        p.text = item['text']
                        p.level = 0  # Bullet level
                        p.font.size = Pt(18)

                        # Enable bullet
                        if hasattr(p, 'bullet'):
                            p.bullet = True

            elif element['type'] == 'block':
                # Handle block environments
                block_content = element.get('content', [])
                if isinstance(block_content, list):
                    # Extract text from content list
                    content_parts = []
                    for item in block_content:
                        if isinstance(item, dict) and 'text' in item:
                            content_parts.append(item['text'])
                        elif isinstance(item, str):
                            content_parts.append(item)
                    content_text = ' '.join(content_parts)
                else:
                    content_text = str(block_content)

                p = text_frame.add_paragraph()
                p.text = content_text
                p.font.size = Pt(16)
                p.font.bold = True

                # Add background shape for block
                self._add_block_background(placeholder)

            elif element['type'] == 'center':
                # Handle centered content
                content_list = element.get('content', [])
                for item in content_list:
                    if isinstance(item, dict) and 'text' in item:
                        text = item['text']
                        # Skip image commands for now
                        if not text.startswith('[width='):
                            p = text_frame.add_paragraph()
                            p.text = text
                            p.alignment = PP_ALIGN.CENTER
                            p.font.size = Pt(18)

    def _add_images_to_slide(self, slide, elements: list):
        """Add images to slide with proper positioning"""
        for element in elements:
            if element['type'] == 'image':
                try:
                    # Try to resolve image path
                    img_path = self._resolve_image_path(element['path'])
                    if img_path:
                        # Add image centered on slide
                        slide.shapes.add_picture(
                            img_path,
                            left=Inches(1),
                            top=Inches(2),
                            width=Inches(8)
                        )
                except Exception as e:
                    logger.warning(f"Could not add image {element['path']}: {e}")

    def _add_block_background(self, placeholder):
        """Add background shape for block environments"""
        # This would require more complex shape manipulation
        # For now, we'll just rely on text formatting
        pass

    def _resolve_image_path(self, path: str) -> str:
        """Resolve image path relative to LaTeX file"""
        import os
        from pathlib import Path

        # Common image directories to check
        tex_dir = Path("latex")
        image_dirs = ["", "images", "..", "../images"]

        for img_dir in image_dirs:
            full_path = tex_dir / img_dir / path
            if full_path.exists():
                return str(full_path)

        return None

def main():
    converter = EnhancedPPTXConverter()
    converter.convert_latex_to_pptx(
        "latex/presentation.tex",
        "latex/enhanced_presentation.pptx"
    )

if __name__ == "__main__":
    main()
