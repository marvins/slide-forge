# Copyright (c) 2026 Slide Forge Team
#
# Permission is hereby granted, free of charge, to any person obtaining a copy
# of this software and associated documentation files (the "Software"), to deal
# in the Software without restriction, including without limitation the rights
# to use, copy, modify, merge, publish, distribute, sublicense, and/or sell
# copies of the Software, and to permit persons to whom the Software is
# furnished to do so, subject to the following conditions:
#
# The above copyright notice and this permission notice shall be included in all
# copies or substantial portions of the Software.
#
# THE SOFTWARE IS PROVIDED "AS IS", WITHOUT WARRANTY OF ANY KIND, EXPRESS OR
# IMPLIED, INCLUDING BUT NOT LIMITED TO THE WARRANTIES OF MERCHANTABILITY,
# FITNESS FOR A PARTICULAR PURPOSE AND NONINFRINGEMENT. IN NO EVENT SHALL THE
# AUTHORS OR COPYRIGHT HOLDERS BE LIABLE FOR ANY CLAIM, DAMAGES OR OTHER
# LIABILITY, WHETHER IN AN ACTION OF CONTRACT, TORT OR OTHERWISE, ARISING FROM,
# OUT OF OR IN CONNECTION WITH THE SOFTWARE OR THE USE OR OTHER DEALINGS IN THE
# SOFTWARE.

"""Tests for PowerPoint builder."""

import pytest
from pathlib import Path
from unittest.mock import Mock, patch
from pptx import Presentation
from slideforge.builders.powerpoint_builder import PowerPoint_Builder
from slideforge.models.universal import (
    Universal_Element, Universal_Frame, Element_Type, Layout_Type, Position, Size
)


class TestPowerPointBuilder:
    """Test cases for PowerPoint builder."""

    @pytest.fixture
    def builder(self):
        """Create PowerPoint builder instance."""
        return PowerPoint_Builder()

    @pytest.fixture
    def output_file(self, tmp_path):
        """Create temporary output file path."""
        return tmp_path / "test_output.pptx"

    @pytest.fixture
    def sample_text_element(self):
        """Create sample text element."""
        return Universal_Element(
            element_type=Element_Type.TEXT,
            content="Sample text content",
            position=Position(x=1.0, y=2.0, width=8.0, height=1.0)
        )

    @pytest.fixture
    def sample_title_element(self):
        """Create sample title element."""
        return Universal_Element(
            element_type=Element_Type.TITLE,
            content="Sample Title"
        )

    @pytest.fixture
    def sample_itemize_element(self):
        """Create sample itemize element."""
        return Universal_Element(
            element_type=Element_Type.ITEMIZE,
            content={'items': ['First item', 'Second item', 'Third item']},
            level=0
        )

    @pytest.fixture
    def sample_image_element(self):
        """Create sample image element."""
        return Universal_Element(
            element_type=Element_Type.IMAGE,
            content={'path': 'test_image.png'},
            position=Position(x=1.0, y=2.0, width=6.0, height=4.0)
        )

    @pytest.fixture
    def sample_equation_element(self):
        """Create sample equation element."""
        return Universal_Element(
            element_type=Element_Type.EQUATION,
            content={'latex': 'E = mc^2', 'type': 'inline'},
            position=Position(x=1.0, y=2.0, width=2.0, height=0.5)
        )

    def test_get_supported_extensions(self, builder):
        """Test supported output extensions."""
        extensions = builder.get_supported_extensions()
        assert '.pptx' in extensions

    def test_get_supported_themes(self, builder):
        """Test supported themes."""
        themes = builder.supported_themes
        assert 'default' in themes
        assert 'professional' in themes
        assert 'academic' in themes
        assert 'minimal' in themes

    def test_get_default_theme(self, builder):
        """Test default theme."""
        assert builder.get_default_theme() == 'default'

    def test_build_empty_presentation(self, builder, output_file):
        """Test building empty presentation."""
        slides = []

        success = builder.build_presentation(slides, output_file, verbose=False)

        assert success
        assert output_file.exists()

        # Verify PowerPoint file structure
        prs = Presentation(str(output_file))
        assert len(prs.slides) == 0

    def test_build_presentation_with_text(self, builder, output_file, sample_text_element):
        """Test building presentation with text elements."""
        from slideforge.models.universal import Universal_Frame

        frame = Universal_Frame(
            frame_number=1,
            title="Test Slide",
            elements=[sample_text_element],
            layout=Layout_Type.TITLE_AND_CONTENT
        )

        success = builder.build_presentation([frame], output_file, verbose=False)

        assert success
        assert output_file.exists()

        # Verify slide was created
        prs = Presentation(str(output_file))
        assert len(prs.slides) == 1

        slide = prs.slides[0]
        assert slide.shapes.title.text == "Test Slide"

    def test_build_presentation_with_itemize(self, builder, output_file, sample_itemize_element):
        """Test building presentation with itemize elements."""
        from slideforge.models.universal import Universal_Frame

        frame = Universal_Frame(
            frame_number=1,
            title="List Slide",
            elements=[sample_itemize_element],
            layout=Layout_Type.TITLE_AND_CONTENT
        )

        success = builder.build_presentation([frame], output_file, verbose=False)

        assert success
        assert output_file.exists()

    def test_build_presentation_with_image(self, builder, output_file, sample_image_element):
        """Test building presentation with image elements."""
        from slideforge.models.universal import Universal_Frame

        frame = Universal_Frame(
            frame_number=1,
            title="Image Slide",
            elements=[sample_image_element],
            layout=Layout_Type.TITLE_AND_CONTENT
        )

        # Mock the image existence check
        with patch('pathlib.Path.exists', return_value=True):
            success = builder.build_presentation([frame], output_file, verbose=False)

        assert success
        assert output_file.exists()

    def test_build_presentation_with_equation(self, builder, output_file, sample_equation_element):
        """Test building presentation with equation elements."""
        from slideforge.models.universal import Universal_Frame

        frame = Universal_Frame(
            frame_number=1,
            title="Equation Slide",
            elements=[sample_equation_element],
            layout=Layout_Type.TITLE_AND_CONTENT
        )

        # Mock the equation rendering
        with patch.object(builder, '_render_latex_equation') as mock_render:
            mock_render.return_value = Path('/fake/path/equation.png')
            with patch('pathlib.Path.exists', return_value=True):
                success = builder.build_presentation([frame], output_file, verbose=False)

        assert success
        assert output_file.exists()
        mock_render.assert_called_once_with('E = mc^2', 'inline', '')

    def test_build_presentation_multiple_slides(self, builder, output_file):
        """Test building presentation with multiple slides."""
        from slideforge.models.universal import Universal_Frame

        slides = []
        for i in range(3):
            frame = Universal_Frame(
                frame_number=i + 1,
                title=f"Slide {i + 1}",
                elements=[Universal_Element(
                    element_type=Element_Type.TEXT,
                    content=f"Content for slide {i + 1}"
                )],
                layout=Layout_Type.TITLE_AND_CONTENT
            )
            slides.append(frame)

        success = builder.build_presentation(slides, output_file, verbose=False)

        assert success
        assert output_file.exists()

        # Verify all slides were created
        prs = Presentation(str(output_file))
        assert len(prs.slides) == 3

        for i, slide in enumerate(prs.slides):
            assert slide.shapes.title.text == f"Slide {i + 1}"

    def test_determine_layout(self, builder):
        """Test layout determination."""
        prs = Presentation()

        # Test title slide
        title_frame = Universal_Frame(
            frame_number=1,
            layout=Layout_Type.TITLE_SLIDE
        )
        layout = builder._determine_layout(title_frame, prs)
        assert layout is not None

        # Test title and content
        content_frame = Universal_Frame(
            frame_number=1,
            layout=Layout_Type.TITLE_AND_CONTENT
        )
        layout = builder._determine_layout(content_frame, prs)
        assert layout is not None

        # Test two column
        two_col_frame = Universal_Frame(
            frame_number=1,
            layout=Layout_Type.TWO_COLUMN
        )
        layout = builder._determine_layout(two_col_frame, prs)
        assert layout is not None

    def test_theme_configuration(self, builder):
        """Test theme configuration."""
        config = builder.theme_configs['default']

        assert 'slide_width' in config
        assert 'slide_height' in config
        assert 'title_font_size' in config
        assert 'content_font_size' in config
        assert 'title_color' in config
        assert 'content_color' in config

        # Test professional theme
        prof_config = builder.theme_configs['professional']
        assert prof_config['title_font_size'] != config['title_font_size']

    def test_invalid_theme(self, builder, output_file):
        """Test handling of invalid theme."""
        from slideforge.models.universal import Universal_Frame

        frame = Universal_Frame(
            frame_number=1,
            title="Test",
            elements=[]
        )

        # Should raise BuilderError for invalid theme
        with pytest.raises(Exception):
            builder.build_presentation([frame], output_file, theme='invalid_theme')

    def test_equation_rendering_latex_not_found(self, builder):
        """Test equation rendering when LaTeX is not available."""
        with patch('subprocess.run') as mock_run:
            # Simulate LaTeX not found
            mock_run.side_effect = FileNotFoundError("latex: command not found")

            result = builder._render_latex_equation('E = mc^2', 'inline', '')
            assert result is None

    def test_equation_rendering_dvipng_not_found(self, builder):
        """Test equation rendering when dvipng is not available."""
        with patch('subprocess.run') as mock_run:
            # First call succeeds (latex), second fails (dvipng)
            mock_run.side_effect = [
                None,  # latex succeeds
                FileNotFoundError("dvipng: command not found")
            ]

            result = builder._render_latex_equation('E = mc^2', 'inline', '')
            assert result is None

    def test_equation_caching(self, builder):
        """Test equation caching functionality."""
        import hashlib
        from pathlib import Path

        latex_eq = "E = mc^2"
        eq_type = "inline"

        with patch('pathlib.Path.exists') as mock_exists:
            # First call: file doesn't exist
            # Second call: file exists (cached)
            mock_exists.side_effect = [False, True]

            with patch('subprocess.run') as mock_run:
                mock_run.return_value.returncode = 0

                with patch('hashlib.md5') as mock_hash:
                    mock_hash_instance = Mock()
                    mock_hash_instance.hexdigest.return_value = 'test_hash'
                    mock_hash.return_value = mock_hash_instance

                    # Call equation rendering
                    result = builder._render_latex_equation(latex_eq, eq_type, '/tmp/test_path')

                    # Verify dvipng was called with relative paths (not full paths)
                    dvipng_calls = [call for call in mock_run.call_args_list
                                   if 'dvipng' in str(call)]
                    if dvipng_calls:
                        dvipng_call = dvipng_calls[0]
                        args = dvipng_call[0][0]  # First argument is the command list
                        assert 'dvipng' in args
                        # Check that relative paths are used (no directory separators in filenames)
                        assert any(arg.endswith('.dvi') and '/' not in arg for arg in args)
                        assert any(arg.endswith('.png') and '/' not in arg for arg in args)

                    result1 = builder._render_latex_equation(latex_eq, eq_type, '/tmp/test_path')
                    result2 = builder._render_latex_equation(latex_eq, eq_type, '/tmp/test_path')

                    # Should return cached result on second call
                    assert result1 == result2
                    assert mock_run.call_count == 1  # Only called once due to caching
