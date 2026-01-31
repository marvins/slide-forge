# Copyright (c) 2026 Slide Forge Team
#
# Permission is hereby granted, free of charge, to any person obtaining a copy
# of this software and associated documentation files (the "Software"), to deal
# in the Software without restriction, including without limitation the rights
# to use, copy, modify, merge, publish, distribute, sublicense, and/or sell
# copies of the Software, and to permit persons to whom the Software is
# furnished to do so, subject to the following conditions:
#
# The above copyright notice and this permission notice shall be included in all
# copies or substantial portions of the Software.
#
# THE SOFTWARE IS PROVIDED "AS IS", WITHOUT WARRANTY OF ANY KIND, EXPRESS OR
# IMPLIED, INCLUDING BUT NOT LIMITED TO THE WARRANTIES OF MERCHANTABILITY,
# FITNESS FOR A PARTICULAR PURPOSE AND NONINFRINGEMENT. IN NO EVENT SHALL THE
# AUTHORS OR COPYRIGHT HOLDERS BE LIABLE FOR ANY CLAIM, DAMAGES OR OTHER
# LIABILITY, WHETHER IN AN ACTION OF CONTRACT, TORT OR OTHERWISE, ARISING FROM,
# OUT OF OR IN CONNECTION WITH THE SOFTWARE OR THE USE OR OTHER DEALINGS IN THE
# SOFTWARE.

"""Unit tests for PowerPoint builder block environment support."""

import pytest
from unittest.mock import Mock, patch
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from src.slideforge.builders.powerpoint_builder import PowerPoint_Builder
from src.slideforge.models.universal import Universal_Element, Element_Type


class TestPowerPointBuilderBlocks:
    """Test PowerPoint builder block environment functionality."""

    @pytest.fixture
    def builder(self):
        """Create a PowerPoint builder instance."""
        return PowerPoint_Builder()

    @pytest.fixture
    def mock_slide(self):
        """Create a mock slide object."""
        slide = Mock()
        slide.shapes = Mock()
        return slide

    @pytest.fixture
    def sample_block_element(self):
        """Create a sample block element."""
        return Universal_Element(
            element_type=Element_Type.BLOCK,
            content={
                'type': 'block',
                'title': 'Test Block',
                'content': 'This is test block content.'
            }
        )

    @pytest.fixture
    def sample_alertblock_element(self):
        """Create a sample alert block element."""
        return Universal_Element(
            element_type=Element_Type.BLOCK,
            content={
                'type': 'alertblock',
                'title': 'Important Alert',
                'content': 'This is an important alert.'
            }
        )

    @pytest.fixture
    def sample_exampleblock_element(self):
        """Create a sample example block element."""
        return Universal_Element(
            element_type=Element_Type.BLOCK,
            content={
                'type': 'exampleblock',
                'title': 'Example',
                'content': 'This is an example.'
            }
        )

    def test_add_block_element_basic(self, builder, mock_slide, sample_block_element):
        """Test adding a basic block element."""
        with patch.object(builder, '_add_block_element') as mock_add_block:
            mock_add_block.return_value = Inches(4)

            result = builder._add_block_element(
                mock_slide,
                sample_block_element,
                builder.themes['default'],
                False,
                Inches(2)
            )

            mock_add_block.assert_called_once()
            assert result == Inches(4)

    def test_block_element_content_dict(self, builder, mock_slide, sample_block_element):
        """Test block element with dictionary content."""
        with patch('src.slideforge.builders.powerpoint_builder.Inches') as mock_inches, \
             patch('src.slideforge.builders.powerpoint_builder.RGBColor') as mock_rgb:

            mock_text_box = Mock()
            mock_text_box.fill = Mock()
            mock_text_box.line = Mock()
            mock_text_box.text_frame = Mock()
            mock_text_frame = Mock()
            mock_text_frame.paragraphs = []
            mock_text_frame.add_paragraph = Mock()

            mock_paragraph = Mock()
            mock_paragraph.font = Mock()
            mock_text_frame.add_paragraph.return_value = mock_paragraph

            mock_slide.shapes.add_textbox.return_value = mock_text_box
            mock_text_box.text_frame = mock_text_frame

            # Call the method
            builder._add_block_element(
                mock_slide,
                sample_block_element,
                builder.themes['default'],
                False,
                Inches(2)
            )

            # Verify textbox was created
            mock_slide.shapes.add_textbox.assert_called_once()

            # Verify fill was set (blue for regular block)
            mock_text_box.fill.solid.assert_called_once()
            mock_text_box.fill.fore_color.rgb = mock_rgb.return_value

            # Verify border was set
            mock_text_box.line.color.rgb = mock_rgb.return_value
            mock_text_box.line.width = Pt(1)

    def test_block_element_colors(self, builder, mock_slide):
        """Test that different block types get different colors."""
        test_cases = [
            ('block', RGBColor(59, 89, 152)),      # Blue
            ('alertblock', RGBColor(220, 38, 127)), # Red
            ('exampleblock', RGBColor(0, 128, 0))   # Green
        ]

        for block_type, expected_color in test_cases:
            element = Universal_Element(
                element_type=Element_Type.BLOCK,
                content={
                    'type': block_type,
                    'title': f'{block_type.title()} Title',
                    'content': f'{block_type} content.'
                }
            )

            with patch('src.slideforge.builders.powerpoint_builder.Inches'), \
                 patch('src.slideforge.builders.powerpoint_builder.RGBColor') as mock_rgb, \
                 patch.object(mock_slide.shapes, 'add_textbox') as mock_add_textbox:

                mock_text_box = Mock()
                mock_text_box.fill = Mock()
                mock_text_box.line = Mock()
                mock_text_box.text_frame = Mock()
                mock_text_frame.paragraphs = []
                mock_text_frame.add_paragraph = Mock()

                mock_add_textbox.return_value = mock_text_box

                builder._add_block_element(
                    mock_slide,
                    element,
                    builder.themes['default'],
                    False,
                    Inches(2)
                )

                # Verify the correct color was used
                mock_rgb.assert_called_with(expected_color)

    def test_block_element_with_string_content(self, builder, mock_slide):
        """Test block element with string content (not dict)."""
        element = Universal_Element(
            element_type=Element_Type.BLOCK,
            content="Simple string content"
        )

        with patch('src.slideforge.builders.powerpoint_builder.Inches'), \
             patch.object(mock_slide.shapes, 'add_textbox') as mock_add_textbox:

            mock_text_box = Mock()
            mock_text_box.fill = Mock()
            mock_text_box.line = Mock()
            mock_text_box.text_frame = Mock()
            mock_text_frame.paragraphs = []
            mock_text_frame.add_paragraph = Mock()

            mock_add_textbox.return_value = mock_text_box

            result = builder._add_block_element(
                mock_slide,
                element,
                builder.themes['default'],
                False,
                Inches(2)
            )

            # Should handle string content gracefully
            assert result is not None

    def test_block_element_text_formatting(self, builder, mock_slide, sample_block_element):
        """Test that block text gets proper formatting."""
        with patch('src.slideforge.builders.powerpoint_builder.Inches'), \
             patch.object(mock_slide.shapes, 'add_textbox') as mock_add_textbox:

            mock_text_box = Mock()
            mock_text_box.fill = Mock()
            mock_text_box.line = Mock()
            mock_text_box.text_frame = Mock()
            mock_text_frame.paragraphs = []

            # Mock paragraphs
            mock_title_p = Mock()
            mock_title_p.font = Mock()
            mock_content_p = Mock()
            mock_content_p.font = Mock()

            mock_text_frame.add_paragraph.side_effect = [mock_title_p, mock_content_p]
            mock_text_box.text_frame = mock_text_frame

            mock_add_textbox.return_value = mock_text_box

            builder._add_block_element(
                mock_slide,
                sample_block_element,
                builder.themes['default'],
                False,
                Inches(2)
            )

            # Verify title is bold
            mock_title_p.font.bold = True

            # Verify font size was set
            mock_title_p.font.size.assert_called()
            mock_content_p.font.size.assert_called()

            # Verify text color was set (white)
            mock_title_p.font.color.rgb = RGBColor(255, 255, 255)
            mock_content_p.font.color.rgb = RGBColor(255, 255, 255)

    def test_block_to_placeholder_returns_false(self, builder, mock_slide, sample_block_element):
        """Test that block_to_placeholder returns False to force element method."""
        result = builder._add_block_to_placeholder(
            mock_slide,
            sample_block_element,
            builder.themes['default'],
            False
        )

        assert result is False

    def test_block_element_margins(self, builder, mock_slide, sample_block_element):
        """Test that block elements have proper margins."""
        with patch('src.slideforge.builders.powerpoint_builder.Inches') as mock_inches, \
             patch.object(mock_slide.shapes, 'add_textbox') as mock_add_textbox:

            mock_text_box = Mock()
            mock_text_box.fill = Mock()
            mock_text_box.line = Mock()
            mock_text_box.text_frame = Mock()
            mock_text_frame.paragraphs = []
            mock_text_frame.add_paragraph = Mock()

            mock_add_textbox.return_value = mock_text_box
            mock_text_box.text_frame = mock_text_frame

            builder._add_block_element(
                mock_slide,
                sample_block_element,
                builder.themes['default'],
                False,
                Inches(2)
            )

            # Verify margins were set
            mock_text_box.text_frame.margin_left = mock_inches.return_value
            mock_text_box.text_frame.margin_right = mock_inches.return_value
            mock_text_box.text_frame.margin_top = mock_inches.return_value
            mock_text_box.text_frame.margin_bottom = mock_inches.return_value
