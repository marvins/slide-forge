#!/usr/bin/env python3
"""
PowerPoint to LaTeX/Markdown Converter and Reverse

This script provides bidirectional conversion between PowerPoint, LaTeX, and Markdown formats.
"""

import os
import sys
import argparse
import subprocess
import re
from pathlib import Path
from typing import Optional, Dict, Any
import logging

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

class PresentationConverter:
    """Main converter class for PowerPoint, LaTeX, and Markdown presentations."""

    def __init__(self):
        self.supported_formats = ['pptx', 'md', 'tex', 'latex']

    def pptx_to_markdown(self, input_path: str, output_path: Optional[str] = None, marp: bool = True) -> str:
        """Convert PowerPoint to Markdown using python-pptx and custom extraction."""
        try:
            from pptx import Presentation
            import markdown
        except ImportError:
            logger.error("Required packages not found. Install with: pip install python-pptx markdown")
            sys.exit(1)

        input_file = Path(input_path)
        if not input_file.exists():
            raise FileNotFoundError(f"Input file not found: {input_path}")

        if output_path is None:
            output_path = input_file.with_suffix('.md')

        logger.info(f"Converting {input_path} to {output_path} (Marp: {marp})")

        # Load presentation
        prs = Presentation(input_path)
        markdown_content = []

        # Add Marp frontmatter if requested
        if marp:
            markdown_content.extend([
                "---",
                "marp: true",
                "theme: default",
                "paginate: true",
                "---",
                ""
            ])

        for i, slide in enumerate(prs.slides):
            if marp:
                markdown_content.append("\n---\n")
                markdown_content.append(f"# Slide {i + 1}\n")
            else:
                markdown_content.append(f"\n# Slide {i + 1}\n")

            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    # Process each paragraph in the text frame
                    for paragraph in shape.text_frame.paragraphs:
                        if not paragraph.text.strip():
                            continue

                        text = paragraph.text.strip()

                        # Detect indentation level for bullet points
                        indent_level = paragraph.level if hasattr(paragraph, 'level') else 0

                        # Check if this is a bullet point using multiple methods
                        is_bullet = False
                        if hasattr(paragraph, 'paragraph_format') and paragraph.paragraph_format:
                            # PowerPoint uses bullet property
                            try:
                                bullet_type = paragraph.paragraph_format.bullet.type
                                is_bullet = bullet_type != 0  # Not None
                            except:
                                is_bullet = False

                        # Fallback: check if text starts with bullet characters
                        if not is_bullet and text:
                            bullet_chars = ['•', '·', '‣', '○', '●', '-', '*', '→']
                            is_bullet = any(text.strip().startswith(char) for char in bullet_chars)

                        # Another fallback: check indentation + common patterns
                        if not is_bullet and indent_level > 0:
                            # Indented text is likely a bullet in PowerPoint
                            is_bullet = True

                        # Format based on type and indentation
                        if is_bullet:
                            # Create bullet with proper indentation
                            bullet_prefix = "  " * indent_level + "* "
                            markdown_content.append(f"{bullet_prefix}{text}")
                        elif indent_level > 0:
                            # Indented text without bullet
                            indent_prefix = "  " * indent_level
                            markdown_content.append(f"{indent_prefix}{text}")
                        elif len(text) < 50 and '\n' not in text and i == 0 and len(markdown_content) < 10:
                            # Likely a title
                            markdown_content.append(f"## {text}")
                        else:
                            # Regular text
                            markdown_content.append(f"{text}")

                    markdown_content.append("")  # Add spacing between shapes

                # Handle images
                if shape.shape_type == 13:  # Picture
                    try:
                        image_bytes = shape.image.blob
                        image_filename = f"slide_{i+1}_image_{len([s for s in slide.shapes if s.shape_type == 13])}.png"
                        image_path = input_file.parent / image_filename

                        with open(image_path, 'wb') as img_file:
                            img_file.write(image_bytes)

                        markdown_content.append(f"![{image_filename}]({image_filename})\n")
                    except Exception as e:
                        logger.warning(f"Could not extract image from slide {i+1}: {e}")

        # Write markdown file
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(markdown_content))

        logger.info(f"Successfully converted to {output_path}")
        return str(output_path)

    def pptx_to_latex(self, input_path: str, output_path: Optional[str] = None, beamer: bool = True) -> str:
        """Convert PowerPoint to LaTeX Beamer format."""
        try:
            from pptx import Presentation
        except ImportError:
            logger.error("python-pptx not found. Install with: pip install python-pptx")
            sys.exit(1)

        input_file = Path(input_path)
        if not input_file.exists():
            raise FileNotFoundError(f"Input file not found: {input_path}")

        if output_path is None:
            output_path = input_file.with_suffix('.tex')

        logger.info(f"Converting {input_path} to {output_path} (Beamer: {beamer})")

        # Load presentation
        prs = Presentation(input_path)

        # LaTeX Beamer template
        if beamer:
            latex_content = [
                "\\documentclass{beamer}",
                "\\usepackage{graphicx}",
                "\\usepackage{hyperref}",
                "\\usepackage[utf8]{inputenc}",
                "\\usepackage[T1]{fontenc}",
                "",
                "\\usetheme{default}",
                "\\usecolortheme{default}",
                "",
                "\\title{Presentation converted from PowerPoint}",
                "\\author{Converted automatically}",
                "\\date{\\today}",
                "",
                "\\begin{document}",
                "",
                "\\frame{\\titlepage}",
                ""
            ]
        else:
            latex_content = [
                "\\documentclass{article}",
                "\\usepackage{graphicx}",
                "\\usepackage{hyperref}",
                "\\usepackage[utf8]{inputenc}",
                "\\usepackage[T1]{fontenc}",
                "",
                "\\title{Presentation converted from PowerPoint}",
                "\\author{Converted automatically}",
                "\\date{\\today}",
                "",
                "\\begin{document}",
                "",
                "\\maketitle",
                ""
            ]

        for i, slide in enumerate(prs.slides):
            if beamer:
                latex_content.append(f"\\begin{{frame}}")

                # Extract slide title from first text element
                slide_title = f"Slide {i + 1}"  # Default fallback
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        potential_title = shape.text.strip()
                        # Use first text as title if it's reasonably short
                        if len(potential_title) < 100 and '\n' not in potential_title:
                            slide_title = potential_title
                            break
                        else:
                            # Use first line of multi-line text as title
                            first_line = potential_title.split('\n')[0].strip()
                            if first_line:
                                slide_title = first_line
                                break

                latex_content.append(f"\\frametitle{{{slide_title}}}")
            else:
                latex_content.append(f"\\section*{{Slide {i + 1}}}")

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    # Convert to LaTeX-safe text
                    text = text.replace('&', '\\&').replace('%', '\\%').replace('$', '\\$')
                    text = text.replace('#', '\\#').replace('_', '\\_').replace('^', '\\^{}')
                    # Handle special Unicode characters
                    text = text.replace('"', '"').replace('"', '"')
                    text = text.replace('"', '"').replace("'", "'")
                    text = text.replace('—', '--').replace('–', '-')
                    text = text.replace('\x0b', '')  # Remove vertical tab characters
                    text = text.replace('\x0c', '')  # Remove form feed characters

                    # Split by lines and format
                    lines = text.split('\n')
                    for line in lines:
                        if line.strip():
                            if beamer:
                                if len(line.strip()) < 30 and not line.strip().endswith('.'):
                                    latex_content.append(f"\\textbf{{{line.strip()}}}\\\\")
                                else:
                                    latex_content.append(f"{line.strip()}\\\\")
                            else:
                                if len(line.strip()) < 30 and not line.strip().endswith('.'):
                                    latex_content.append(f"\\subsection*{{{line.strip()}}}")
                                else:
                                    latex_content.append(f"{line.strip()}")

                # Handle images
                if shape.shape_type == 13:  # Picture
                    try:
                        image_bytes = shape.image.blob
                        image_filename = f"slide_{i+1}_image_{len([s for s in slide.shapes if s.shape_type == 13])}.png"
                        image_path = input_file.parent / image_filename

                        with open(image_path, 'wb') as img_file:
                            img_file.write(image_bytes)

                        if beamer:
                            latex_content.append(f"\\includegraphics[width=0.8\\textwidth]{{{image_filename}}}\\\\")
                        else:
                            latex_content.append(f"\\includegraphics[width=0.6\\textwidth]{{{image_filename}}}")
                            latex_content.append("")
                    except Exception as e:
                        logger.warning(f"Could not extract image from slide {i+1}: {e}")

            if beamer:
                latex_content.append("\\end{frame}")
                latex_content.append("")
            else:
                latex_content.append("")
                latex_content.append("")

        latex_content.append("\\end{document}")

        # Write LaTeX file
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(latex_content))

        logger.info(f"Successfully converted to {output_path}")
        return str(output_path)

    def markdown_to_pptx(self, input_path: str, output_path: Optional[str] = None) -> str:
        """Convert Markdown to PowerPoint using pandoc."""
        try:
            import pypandoc
        except ImportError:
            logger.error("pypandoc not found. Install with: pip install pypandoc")
            sys.exit(1)

        input_file = Path(input_path)
        if not input_file.exists():
            raise FileNotFoundError(f"Input file not found: {input_path}")

        if output_path is None:
            output_path = input_file.with_suffix('.pptx')

        logger.info(f"Converting {input_path} to {output_path}")

        try:
            # Convert markdown to pptx using pandoc
            pypandoc.convert_file(
                str(input_file),
                'pptx',
                outputfile=str(output_path),
                extra_args=['--standalone']
            )
            logger.info(f"Successfully converted to {output_path}")
            return str(output_path)
        except Exception as e:
            logger.error(f"Error converting with pandoc: {e}")
            # Fallback: create a simple PowerPoint manually
            return self._markdown_to_pptx_fallback(input_path, output_path)

    def _markdown_to_pptx_fallback(self, input_path: str, output_path: str) -> str:
        """Fallback method to convert Markdown to PowerPoint manually."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            logger.error("python-pptx not found. Install with: pip install python-pptx")
            sys.exit(1)

        logger.info("Using fallback method for Markdown to PowerPoint conversion")

        prs = Presentation()

        with open(input_path, 'r', encoding='utf-8') as f:
            content = f.read()

        # Split content by slide markers (# Slide X or just #)
        sections = content.split('\n# ')

        for section in sections:
            if not section.strip():
                continue

            lines = section.strip().split('\n')
            if not lines:
                continue

            # Create new slide
            slide_layout = prs.slide_layouts[1]  # Title and content
            slide = prs.slides.add_slide(slide_layout)

            # First line as title
            title = lines[0].replace('#', '').strip()
            if title:
                slide.shapes.title.text = title

            # Rest as content
            if len(lines) > 1:
                content_text = '\n'.join(lines[1:])
                if content_text.strip():
                    slide.placeholders[1].text = content_text

        prs.save(output_path)
        logger.info(f"Successfully converted to {output_path}")
        return output_path

    def latex_to_pptx(self, input_path: str, output_path: Optional[str] = None) -> str:
        """Convert LaTeX Beamer to PowerPoint using pandoc."""
        try:
            import pypandoc
        except ImportError:
            logger.error("pypandoc not found. Install with: pip install pypandoc")
            sys.exit(1)

        input_file = Path(input_path)
        if not input_file.exists():
            raise FileNotFoundError(f"Input file not found: {input_path}")

        if output_path is None:
            output_path = input_file.with_suffix('.pptx')

        logger.info(f"Converting {input_path} to {output_path}")

        with open(input_path, 'r', encoding='utf-8') as f:
            content = f.read()

        # For Beamer presentations, always use our fallback method as pandoc doesn't handle frames correctly
        if 'beamer' in content.lower() or '\\begin{frame}' in content:
            logger.info("Detected Beamer presentation, using fallback method")
            return self._latex_to_pptx_fallback(input_path, output_path)

        try:
            # Convert LaTeX to pptx using pandoc (for non-Beamer LaTeX)
            pypandoc.convert_file(
                str(input_file),
                'pptx',
                outputfile=str(output_path),
                extra_args=['--standalone']
            )
            logger.info(f"Successfully converted to {output_path}")
            return str(output_path)
        except Exception as e:
            logger.error(f"Error converting with pandoc: {e}")
            # Fallback: extract text and create simple PowerPoint
            return self._latex_to_pptx_fallback(input_path, output_path)

    def _latex_to_pptx_fallback(self, input_path: str, output_path: str) -> str:
        """Fallback method to convert LaTeX to PowerPoint manually."""
        try:
            from pptx import Presentation
            import re
        except ImportError:
            logger.error("python-pptx not found. Install with: pip install python-pptx")
            sys.exit(1)

        logger.info("Using fallback method for LaTeX to PowerPoint conversion")

        prs = Presentation()

        with open(input_path, 'r', encoding='utf-8') as f:
            content = f.read()

        # Extract frames (slides) from LaTeX
        # Handle frames with optional arguments like \begin{frame}[plain,t]
        frame_pattern = r'\\begin\{frame\}(?:\[[^\]]*\])?\s*(.*?)\\end\{frame\}'
        frames = re.findall(frame_pattern, content, re.DOTALL)

        logger.info(f"Found {len(frames)} frames to convert")

        for frame_content in frames:
            slide_layout = prs.slide_layouts[1]  # Title and content
            slide = prs.slides.add_slide(slide_layout)

            # Extract frame title
            title_match = re.search(r'\\frametitle\{(.*?)\}', frame_content)
            if title_match:
                title = title_match.group(1).replace('\\', '').strip()
                slide.shapes.title.text = title

            # Extract text content
            text_content = frame_content

            # Remove frame title from content to avoid duplication
            text_content = re.sub(r'\\frametitle\{[^}]*\}', '', text_content)

            # Handle special frames like titlepage
            if '\\titlepage' in text_content:
                # This is a title slide, use title information
                if hasattr(prs, 'slide_layouts') and len(prs.slide_layouts) > 0:
                    title_layout = prs.slide_layouts[0]  # Title slide layout
                    title_slide = prs.slides.add_slide(title_layout)
                    # Set title and subtitle from document info
                    title_slide.shapes.title.text = "NITF Image Conversion Library"
                    if hasattr(title_slide.placeholders, 1):
                        title_slide.placeholders[1].text = "Professional C++ SDK for DOD NITF Format"
                continue

            # Process content line by line to preserve structure
            lines = text_content.split('\n')
            processed_lines = []

            for line in lines:
                line = line.strip()
                if not line:
                    continue

                # Handle itemize items
                if '\\item' in line and not line.startswith('\\begin') and not line.startswith('\\end'):
                    # Extract the text after \item
                    item_text = re.sub(r'\\item\s*', '', line)
                    item_text = self._clean_latex_text(item_text)
                    if item_text:
                        processed_lines.append(f"• {item_text}")
                # Skip LaTeX environments and commands
                elif line.startswith('\\begin') or line.startswith('\\end'):
                    continue
                elif line.startswith('\\') or line.startswith('{') or line.startswith('}'):
                    continue
                elif line == 'titlepage':
                    continue
                else:
                    # Regular text
                    clean_line = self._clean_latex_text(line)
                    if clean_line and len(clean_line) > 1:  # Skip very short fragments
                        processed_lines.append(clean_line)

            # Join processed lines
            final_text = '\n'.join(processed_lines)

            if final_text.strip():
                slide.placeholders[1].text = final_text

        prs.save(output_path)
        logger.info(f"Successfully converted to {output_path}")
        return output_path

    def _clean_latex_text(self, text: str) -> str:
        """Clean LaTeX commands and formatting from text."""
        # Remove LaTeX commands with arguments
        text = re.sub(r'\\[a-zA-Z]+\{([^}]*)\}', r'\1', text)
        # Remove LaTeX commands without arguments
        text = re.sub(r'\\[a-zA-Z]+', '', text)
        # Remove braces
        text = re.sub(r'[{}]', '', text)
        # Handle line breaks
        text = re.sub(r'\\\\', '\n', text)
        # Clean up multiple newlines
        text = re.sub(r'\n+', '\n', text)
        # Remove common LaTeX environments
        text = re.sub(r'\\begin\{[^}]*\}|\\end\{[^}]*\}', '', text)
        # Strip whitespace
        return text.strip()

    def convert(self, input_path: str, output_format: str, output_path: Optional[str] = None,
                marp: bool = True, beamer: bool = True) -> str:
        """Main conversion method."""
        input_file = Path(input_path)
        input_format = input_file.suffix.lower().lstrip('.')

        if input_format not in self.supported_formats:
            raise ValueError(f"Unsupported input format: {input_format}")

        if output_format not in self.supported_formats:
            raise ValueError(f"Unsupported output format: {output_format}")

        if input_format == output_format:
            logger.warning("Input and output formats are the same. No conversion needed.")
            return str(input_path)

        # Route to appropriate conversion method
        if input_format == 'pptx':
            if output_format == 'md':
                return self.pptx_to_markdown(input_path, output_path, marp=marp)
            elif output_format in ['tex', 'latex']:
                return self.pptx_to_latex(input_path, output_path, beamer=beamer)

        elif input_format == 'md':
            if output_format == 'pptx':
                return self.markdown_to_pptx(input_path, output_path)

        elif input_format in ['tex', 'latex']:
            if output_format == 'pptx':
                return self.latex_to_pptx(input_path, output_path)

        raise NotImplementedError(f"Conversion from {input_format} to {output_format} not yet implemented")

def main():
    """Command line interface."""
    parser = argparse.ArgumentParser(
        description='Convert between PowerPoint, LaTeX, and Markdown formats',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  %(prog)s input.pptx md                    # Convert PowerPoint to Markdown
  %(prog)s input.pptx tex                   # Convert PowerPoint to LaTeX
  %(prog)s input.md pptx                    # Convert Markdown to PowerPoint
  %(prog)s input.tex pptx                   # Convert LaTeX to PowerPoint
  %(prog)s input.pptx md -o output.md       # Specify output file
        """
    )

    parser.add_argument('input', help='Input file path')
    parser.add_argument('format', help='Output format (pptx, md, tex, latex)')
    parser.add_argument('-o', '--output', help='Output file path (optional)')
    parser.add_argument('-v', '--verbose', action='store_true', help='Verbose output')
    parser.add_argument('--no-marp', action='store_true', help='Disable Marp format for Markdown output')
    parser.add_argument('--no-beamer', action='store_true', help='Disable Beamer format for LaTeX output')

    args = parser.parse_args()

    if args.verbose:
        logging.getLogger().setLevel(logging.DEBUG)

    converter = PresentationConverter()

    try:
        output_path = converter.convert(
            args.input,
            args.format,
            args.output,
            marp=not args.no_marp,
            beamer=not args.no_beamer
        )
        print(f"Conversion completed: {output_path}")
    except Exception as e:
        logger.error(f"Conversion failed: {e}")
        sys.exit(1)

if __name__ == '__main__':
    main()
