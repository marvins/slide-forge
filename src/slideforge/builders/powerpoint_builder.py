# Copyright (c) 2026 Slide Forge Team
#
# Permission is hereby granted, free of charge, to any person obtaining a copy
# of this software and associated documentation files (the "Software"), to deal
# in the Software without restriction, including without limitation the rights
# to use, copy, modify, merge, publish, distribute, sublicense, and/or sell
# copies of the Software, and to permit persons to whom the Software is
# furnished to do so, subject to the following conditions:
#
# The above copyright notice and this permission notice shall be included in all
# copies or substantial portions of the Software.
#
# THE SOFTWARE IS PROVIDED "AS IS", WITHOUT WARRANTY OF ANY KIND, EXPRESS OR
# IMPLIED, INCLUDING BUT NOT LIMITED TO THE WARRANTIES OF MERCHANTABILITY,
# FITNESS FOR A PARTICULAR PURPOSE AND NONINFRINGEMENT. IN NO EVENT SHALL THE
# AUTHORS OR COPYRIGHT HOLDERS BE LIABLE FOR ANY CLAIM, DAMAGES OR OTHER
# LIABILITY, WHETHER IN AN ACTION OF CONTRACT, TORT OR OTHERWISE, ARISING FROM,
# OUT OF OR IN CONNECTION WITH THE SOFTWARE OR THE USE OR OTHER DEALINGS IN THE
# SOFTWARE.

"""PowerPoint builder implementation using python-pptx."""

from pathlib import Path
from typing import List, Dict, Any, Optional
import logging

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from ..base import Base_Builder
from ..models.universal import Universal_Frame, Universal_Element, Element_Type, Layout_Type
from ..exceptions import BuilderError


class PowerPoint_Builder(Base_Builder):
    """Builder for PowerPoint presentations using python-pptx."""

    def __init__(self):
        """Initialize PowerPoint builder."""
        self.supported_themes = ['default', 'professional', 'academic', 'minimal']
        self.default_theme = 'default'
        self.logger = logging.getLogger(__name__)

        # Theme configurations
        self.theme_configs = {
            'default': {
                'slide_width': Inches(10),
                'slide_height': Inches(7.5),
                'title_font_size': 44,
                'content_font_size': 18,
                'title_color': RGBColor(0, 0, 0),
                'content_color': RGBColor(0, 0, 0),
                'background_color': RGBColor(255, 255, 255)
            },
            'professional': {
                'slide_width': Inches(10),
                'slide_height': Inches(7.5),
                'title_font_size': 40,
                'content_font_size': 16,
                'title_color': RGBColor(0, 32, 96),
                'content_color': RGBColor(32, 32, 32),
                'background_color': RGBColor(255, 255, 255)
            },
            'academic': {
                'slide_width': Inches(10),
                'slide_height': Inches(7.5),
                'title_font_size': 42,
                'content_font_size': 17,
                'title_color': RGBColor(0, 0, 128),
                'content_color': RGBColor(0, 0, 0),
                'background_color': RGBColor(255, 255, 255)
            },
            'minimal': {
                'slide_width': Inches(10),
                'slide_height': Inches(7.5),
                'title_font_size': 36,
                'content_font_size': 14,
                'title_color': RGBColor(64, 64, 64),
                'content_color': RGBColor(64, 64, 64),
                'background_color': RGBColor(255, 255, 255)
            }
        }

    def build_presentation(self, slides: List[Universal_Frame],
                          output_file: Path, **kwargs) -> bool:
        """
        Build a PowerPoint presentation from slide structures.

        Args:
            slides: List of Universal_Frame objects
            output_file: Path to output .pptx file
            **kwargs: Additional build options
                - theme: PowerPoint theme
                - preserve_colors: Preserve colors from source
                - include_images: Include images from source
                - verbose: Enable verbose output

        Returns:
            True if build successful, False otherwise

        Raises:
            BuilderError: If build fails
        """
        try:
            # Get options
            theme = kwargs.get('theme', self.default_theme)
            preserve_colors = kwargs.get('preserve_colors', True)
            include_images = kwargs.get('include_images', True)
            verbose = kwargs.get('verbose', False)
            source_path = kwargs.get('source_path', '')  # Get source path for image resolution

            if verbose:
                self.logger.info(f"Building PowerPoint presentation with theme: {theme}")
                self.logger.info(f"Output file: {output_file}")
                self.logger.info(f"Number of slides: {len(slides)}")
                if source_path:
                    self.logger.info(f"Source path: {source_path}")

            # Validate theme
            if theme not in self.supported_themes:
                raise BuilderError(f"Unsupported theme: {theme}",
                               operation="build_presentation", output_format="pptx")

            # Get theme configuration
            config = self.theme_configs[theme]

            # Create presentation
            prs = Presentation()

            # Set slide dimensions
            prs.slide_width = config['slide_width']
            prs.slide_height = config['slide_height']

            # Build each slide
            for i, slide in enumerate(slides, 1):
                if verbose:
                    self.logger.info(f"Building slide {i}: {slide.title or 'No title'}")

                slide_layout = self._determine_layout(slide, prs)
                slide_obj = prs.slides.add_slide(slide_layout)

                # Add title if present
                if slide.title:
                    title_shape = slide_obj.shapes.title
                    title_shape.text = slide.title
                    if preserve_colors:
                        title_shape.text_frame.paragraphs[0].font.color.rgb = config['title_color']
                    # Set font size with proper conversion
                    font_size = config.get('title_font_size', 44)
                    if font_size > 0:
                        title_shape.text_frame.paragraphs[0].font.size = Pt(font_size)

                # Add elements to slide, using content placeholder when possible
                self._add_elements_to_slide(slide_obj, slide.elements, config, preserve_colors, include_images, source_path)

            # Ensure output directory exists
            output_file.parent.mkdir(parents=True, exist_ok=True)

            # Save presentation
            prs.save(output_file)

            if verbose:
                self.logger.info(f"Successfully built PowerPoint presentation: {output_file}")

            return True

        except Exception as e:
            raise BuilderError(f"Failed to build PowerPoint presentation: {e}",
                           operation="build_presentation", output_format="pptx")

    def get_supported_extensions(self) -> List[str]:
        """Get supported output extensions."""
        return ['.pptx']

    def get_default_theme(self) -> str:
        """Get default theme."""
        return self.default_theme

    def _determine_layout(self, slide: Universal_Frame, presentation):
        """Determine the appropriate slide layout for the frame."""
        if slide.layout == Layout_Type.TITLE_SLIDE:
            return presentation.slide_layouts[0]  # Title Slide
        elif slide.layout == Layout_Type.TITLE_AND_CONTENT:
            return presentation.slide_layouts[1]  # Title and Content
        elif slide.layout == Layout_Type.TWO_COLUMN:
            return presentation.slide_layouts[3]  # Two Content
        else:
            return presentation.slide_layouts[1]  # Title and Content (default)

    def _add_elements_to_slide(self, slide_obj, elements: List[Universal_Element],
                              config: Dict[str, Any], preserve_colors: bool, include_images: bool, source_path: str = ''):
        """Add elements to a PowerPoint slide using native placeholders when possible."""
        content_placeholder_used = False

        for element in elements:
            try:
                if element.element_type == Element_Type.TEXT:
                    # Use content placeholder for first text element if available
                    if not content_placeholder_used and hasattr(slide_obj.shapes, 'placeholders'):
                        content_placeholder_used = self._add_text_to_placeholder(slide_obj, element, config, preserve_colors)
                    else:
                        self._add_text_element(slide_obj, element, config, preserve_colors)
                elif element.element_type == Element_Type.TITLE:
                    # Title elements are handled separately by slide title
                    pass
                elif element.element_type == Element_Type.ITEMIZE:
                    # Use content placeholder for first itemize if available
                    if not content_placeholder_used and hasattr(slide_obj.shapes, 'placeholders'):
                        content_placeholder_used = self._add_itemize_to_placeholder(slide_obj, element, config, preserve_colors)
                    else:
                        self._add_itemize_element(slide_obj, element, config, preserve_colors)
                elif element.element_type == Element_Type.IMAGE and include_images:
                    self._add_image_element(slide_obj, element, config, source_path)
                elif element.element_type == Element_Type.EQUATION and include_images:
                    self._add_equation_element(slide_obj, element, config, source_path)
                elif element.element_type == Element_Type.BLOCK:
                    # Use content placeholder for first block if available
                    if not content_placeholder_used and hasattr(slide_obj.shapes, 'placeholders'):
                        content_placeholder_used = self._add_block_to_placeholder(slide_obj, element, config, preserve_colors)
                    else:
                        # For blocks, we need to track positioning
                        current_top = Inches(2.5)  # Start below title
                        self._add_block_element(slide_obj, element, config, preserve_colors, current_top)
            except Exception as e:
                self.logger.warning(f"Failed to add element {element.element_type}: {e}")

    def _add_text_element(self, slide_obj, element: Universal_Element,
                          config: Dict[str, Any], preserve_colors: bool):
        """Add a text element to the slide using its predefined position."""
        content = element.content
        if isinstance(content, str):
            text = content
        else:
            text = content.text

        # Use position from element if available, otherwise fallback
        if element.position:
            left = Inches(element.position.x)
            top = Inches(element.position.y)
            width = Inches(element.position.width) if element.position.width else Inches(8)
            height = Inches(element.position.height) if element.position.height else Inches(1.5)
        else:
            # Fallback positioning
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(1.5)

        text_box = slide_obj.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True  # Enable text wrapping

        # Split text by paragraphs and add each as a separate paragraph
        paragraphs = text.split('\n\n') if '\n\n' in text else [text]

        for i, para_text in enumerate(paragraphs):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            p.text = para_text.strip()

            # Set font size with proper conversion
            font_size = config.get('content_font_size', 18)
            if font_size > 0:
                p.font.size = Pt(font_size)

            if preserve_colors:
                p.font.color.rgb = config['content_color']

    def _add_title_element(self, slide_obj, element: Universal_Element,
                           config: Dict[str, Any], preserve_colors: bool):
        """Add a title element to the slide."""
        content = element.content
        if isinstance(content, str):
            text = content
        else:
            text = content.text

        # Add to existing title shape if available, otherwise create new one
        if slide_obj.shapes.title:
            title_shape = slide_obj.shapes.title
            title_shape.text = text
            if preserve_colors:
                title_shape.text_frame.paragraphs[0].font.color.rgb = config['title_color']
            # Set font size with proper conversion
            font_size = config.get('title_font_size', 44)
            if font_size > 0:
                title_shape.text_frame.paragraphs[0].font.size = Pt(font_size)
        else:
            # Create new title text box
            left = Inches(1) if element.position else Inches(1)
            top = Inches(0.5) if element.position else Inches(0.5)
            width = Inches(8) if element.size else Inches(8)
            height = Inches(1) if element.size else Inches(1)

            title_box = slide_obj.shapes.add_textbox(left, top, width, height)
            title_frame = title_box.text_frame
            p = title_frame.paragraphs[0]
            p.text = text
            # Set font size with proper conversion
            font_size = config.get('title_font_size', 44)
            if font_size > 0:
                p.font.size = font_size
            p.font.bold = True

            if preserve_colors:
                p.font.color.rgb = config['title_color']

    def _add_itemize_element(self, slide_obj, element: Universal_Element,
                           config: Dict[str, Any], preserve_colors: bool):
        """Add a bullet list element to the slide using its predefined position."""
        content = element.content
        if isinstance(content, dict) and 'items' in content:
            items = content['items']
        else:
            # Try to parse as string
            items = [str(content)]

        # Use position from element if available, otherwise fallback
        if element.position:
            left = Inches(element.position.x)
            top = Inches(element.position.y)
            width = Inches(element.position.width) if element.position.width else Inches(8)
            height = Inches(element.position.height) if element.position.height else Inches(max(0.5, len(items) * 0.4))
        else:
            # Fallback positioning
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(max(0.5, len(items) * 0.4))

        text_box = slide_obj.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame

        for i, item in enumerate(items):
            if i > 0:
                p = text_frame.add_paragraph()
            else:
                p = text_frame.paragraphs[0]

            p.text = item
            p.level = element.level if hasattr(element, 'level') else 0
            # Set font size with proper conversion
            font_size = config.get('content_font_size', 18)
            if font_size > 0:
                p.font.size = Pt(font_size)

            if preserve_colors:
                p.font.color.rgb = config['content_color']

    def _add_image_element(self, slide_obj, element: Universal_Element,
                          config: Dict[str, Any], source_path: str = '', current_top = Inches(2)):
        """Add an image element to the slide and return the new top position."""
        content = element.content
        if isinstance(content, dict) and 'path' in content:
            image_path = content['path']
        else:
            image_path = str(content)

        try:
            # Convert relative paths to absolute
            if not Path(image_path).is_absolute():
                if source_path:
                    # Resolve relative to source document directory
                    source_dir = Path(source_path).parent
                    image_path = source_dir / image_path
                else:
                    # Fallback to current working directory
                    image_path = Path.cwd() / image_path

            if Path(image_path).exists():
                left = Inches(1) if element.position else Inches(1)
                top = current_top
                width = Inches(6) if element.size else Inches(6)
                height = Inches(4) if element.size else None

                slide_obj.shapes.add_picture(str(image_path), left, top, width, height)

                # Return the new top position (below this image)
                image_height = height if height else Inches(4)
                return top + image_height
            else:
                self.logger.warning(f"Image file not found: {image_path}")
                return current_top
        except Exception as e:
            self.logger.warning(f"Failed to add image {image_path}: {e}")
            return current_top

    def _add_block_element(self, slide_obj, element: Universal_Element,
                          config: Dict[str, Any], preserve_colors: bool, current_top):
        """Add a Beamer-style block element to the slide and return the new top position."""
        content = element.content

        # Extract title, content, and type from block
        if isinstance(content, dict):
            block_type = content.get('type', 'block')
            block_title = content.get('title', 'Block')
            block_content = content.get('content', '')
        else:
            block_type = 'block'
            block_title = 'Block'
            block_content = str(content)

        left = Inches(1) if element.position else Inches(1)
        top = current_top
        width = Inches(8) if element.size else Inches(8)
        height = Inches(1.5) if element.size else Inches(1.5)  # Taller for blocks

        # Create text box with Beamer-style formatting
        text_box = slide_obj.shapes.add_textbox(left, top, width, height)

        # Apply Beamer block styling based on type
        fill = text_box.fill
        fill.solid()

        # Set colors based on block type
        if block_type == 'alertblock':
            fill.fore_color.rgb = RGBColor(220, 38, 127)  # Beamer alert red
            text_color = RGBColor(255, 255, 255)  # White text
        elif block_type == 'exampleblock':
            fill.fore_color.rgb = RGBColor(0, 128, 0)  # Beamer example green
            text_color = RGBColor(255, 255, 255)  # White text
        else:  # regular block
            fill.fore_color.rgb = RGBColor(59, 89, 152)  # Beamer blue background
            text_color = RGBColor(255, 255, 255)  # White text

        # Add border
        line = text_box.line
        line.color.rgb = RGBColor(0, 0, 0)  # Black border
        line.width = Pt(1)  # Thin border

        text_frame = text_box.text_frame
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_bottom = Inches(0.1)

        # Add title paragraph (bold, white)
        if block_title:
            title_p = text_frame.add_paragraph()
            title_p.text = block_title
            title_p.font.bold = True
            title_p.font.color.rgb = text_color
            title_font_size = config.get('content_font_size', 18)
            if title_font_size > 0:
                title_p.font.size = Pt(title_font_size)

        # Add content paragraph (white)
        if block_content:
            content_p = text_frame.add_paragraph()
            content_p.text = block_content
            content_p.font.color.rgb = text_color
            content_font_size = config.get('content_font_size', 18)
            if content_font_size > 0:
                content_p.font.size = Pt(content_font_size)

        # Return the new top position (below this element)
        return top + height

    def _add_text_to_placeholder(self, slide_obj, element: Universal_Element,
                                config: Dict[str, Any], preserve_colors: bool) -> bool:
        """Add text to the slide's content placeholder. Returns True if successful."""
        try:
            # Find the content placeholder (check both type and name for robustness)
            for placeholder in slide_obj.placeholders:
                # Check for body/content placeholder by type or name
                is_body_placeholder = (
                    placeholder.placeholder_format.type == 1 or  # Body placeholder
                    'content' in placeholder.name.lower() or
                    'body' in placeholder.name.lower()
                )
                # Skip title placeholder
                is_title_placeholder = (
                    placeholder.placeholder_format.type == 0 or  # Title placeholder
                    'title' in placeholder.name.lower()
                )

                if is_body_placeholder and not is_title_placeholder:
                    content = element.content
                    if isinstance(content, str):
                        text = content
                    else:
                        text = content.text

                    placeholder.text = text

                    # Apply formatting to the placeholder
                    text_frame = placeholder.text_frame
                    for paragraph in text_frame.paragraphs:
                        # Set font size with proper conversion
                        font_size = config.get('content_font_size', 18)
                        if font_size > 0:
                            paragraph.font.size = Pt(font_size)

                        if preserve_colors:
                            paragraph.font.color.rgb = config['content_color']

                    return True
        except Exception as e:
            self.logger.warning(f"Failed to use content placeholder: {e}")

        return False

    def _add_itemize_to_placeholder(self, slide_obj, element: Universal_Element,
                                    config: Dict[str, Any], preserve_colors: bool) -> bool:
        """Add itemize to the slide's content placeholder. Returns True if successful."""
        try:
            # Find the content placeholder (check both type and name for robustness)
            for placeholder in slide_obj.placeholders:
                # Check for body/content placeholder by type or name
                is_body_placeholder = (
                    placeholder.placeholder_format.type == 1 or  # Body placeholder
                    'content' in placeholder.name.lower() or
                    'body' in placeholder.name.lower()
                )
                # Skip title placeholder
                is_title_placeholder = (
                    placeholder.placeholder_format.type == 0 or  # Title placeholder
                    'title' in placeholder.name.lower()
                )

                if is_body_placeholder and not is_title_placeholder:
                    content = element.content
                    if isinstance(content, dict) and 'items' in content:
                        items = content['items']
                    else:
                        items = [str(content)]

                    text_frame = placeholder.text_frame
                    text_frame.clear()  # Clear existing content

                    for i, item in enumerate(items):
                        if i > 0:
                            p = text_frame.add_paragraph()
                        else:
                            p = text_frame.paragraphs[0]

                        p.text = item
                        p.level = element.level if hasattr(element, 'level') else 0

                        # Set font size with proper conversion
                        font_size = config.get('content_font_size', 18)
                        if font_size > 0:
                            p.font.size = Pt(font_size)

                        if preserve_colors:
                            p.font.color.rgb = config['content_color']

                    return True
        except Exception as e:
            self.logger.warning(f"Failed to use content placeholder for itemize: {e}")

        return False

    def _add_block_to_placeholder(self, slide_obj, element: Universal_Element,
                                config: Dict[str, Any], preserve_colors: bool) -> bool:
        """Add block to the slide's content placeholder. Returns True if successful."""
        try:
            # For blocks, we prefer creating a separate text box with styling
            # rather than using the placeholder, to get the visual effect
            return False  # Force fallback to _add_block_element
        except Exception as e:
            self.logger.warning(f"Failed to add block to placeholder: {e}")
            return False

    def _add_equation_element(self, slide_obj, element: Universal_Element,
                            config: Dict[str, Any], source_path: str = ''):
        """Add an equation element by rendering LaTeX to image."""
        try:
            content = element.content
            latex_equation = content.get('latex', '')
            equation_type = content.get('type', 'inline')

            if not latex_equation:
                self.logger.warning("Empty equation content")
                return

            # Render LaTeX equation to image
            image_path = self._render_latex_equation(latex_equation, equation_type, source_path)

            if image_path and Path(image_path).exists():
                # Add as image with equation-specific positioning
                left = Inches(1) if element.position else Inches(1)
                top = Inches(2) if element.position else Inches(2)

                # Different sizing for inline vs display equations
                if equation_type == 'inline':
                    width = Inches(2) if element.size else Inches(2)
                    height = Inches(0.5) if element.size else Inches(0.5)
                else:  # display
                    width = Inches(6) if element.size else Inches(6)
                    height = Inches(1.5) if element.size else Inches(1.5)

                slide_obj.shapes.add_picture(str(image_path), left, top, width, height)
                self.logger.info(f"Successfully added equation: {latex_equation[:50]}...")
            else:
                self.logger.warning(f"Failed to render equation image: {latex_equation[:50]}...")

        except Exception as e:
            self.logger.warning(f"Failed to add equation element: {e}")

    def _render_latex_equation(self, latex_equation: str, equation_type: str, source_path: str = '') -> Optional[Path]:
        """Render LaTeX equation to PNG image using temporary files."""
        import tempfile
        import subprocess
        import hashlib

        try:
            # Create cache directory
            cache_dir = Path(source_path).parent / '.equation_cache' if source_path else Path.cwd() / '.equation_cache'
            cache_dir.mkdir(exist_ok=True)

            # Create hash for caching
            equation_hash = hashlib.md5(f"{latex_equation}_{equation_type}".encode()).hexdigest()
            cached_image = cache_dir / f"eq_{equation_hash}.png"

            # Return cached image if exists
            if cached_image.exists():
                return cached_image

            # Create temporary LaTeX document
            if equation_type == 'inline':
                latex_doc = f"""
\\documentclass[preview]{{standalone}}
\\usepackage{{amsmath}}
\\begin{{document}}
\\({latex_equation}\\)
\\end{{document}}
"""
            else:  # display
                latex_doc = f"""
\\documentclass[preview]{{standalone}}
\\usepackage{{amsmath}}
\\begin{{document}}
\\[{latex_equation}\\]
\\end{{document}}
"""

            # Create temporary files
            with tempfile.NamedTemporaryFile(mode='w', suffix='.tex', delete=False, dir=cache_dir) as f:
                f.write(latex_doc)
                tex_path = Path(f.name)

            # Compile LaTeX to DVI
            result = subprocess.run([
                'latex', '-output-directory=' + str(cache_dir), '-interaction=nonstopmode', str(tex_path)
            ], capture_output=True, text=True, cwd=cache_dir)

            if result.returncode != 0:
                self.logger.warning(f"LaTeX compilation failed: {result.stderr}")
                return None

            # Convert DVI to PNG
            dvi_path = cache_dir / tex_path.stem
            png_path = cache_dir / f"eq_{equation_hash}.png"

            result = subprocess.run([
                'dvipng', '-T', 'tight', '-D', '300', '-bg', 'White',
                '-Q', '9', '-o', png_path.name, dvi_path.name
            ], capture_output=True, text=True, cwd=cache_dir)

            if result.returncode != 0:
                self.logger.warning(f"DVI to PNG conversion failed: {result.stderr}")
                return None

            # Cleanup temporary files
            tex_path.unlink(missing_ok=True)
            dvi_path.with_suffix('.dvi').unlink(missing_ok=True)
            dvi_path.with_suffix('.log').unlink(missing_ok=True)
            dvi_path.with_suffix('.aux').unlink(missing_ok=True)

            return png_path if png_path.exists() else None

        except FileNotFoundError as e:
            if 'latex' in str(e):
                self.logger.warning("LaTeX not found. Equations will be skipped. Install LaTeX (TeX Live, MiKTeX, etc.)")
            elif 'dvipng' in str(e):
                self.logger.warning("dvipng not found. Equation images will be skipped. Install dvipng")
            else:
                self.logger.warning(f"File not found in equation rendering: {e}")
            return None
        except Exception as e:
            self.logger.warning(f"Error rendering equation: {e}")
            return None
